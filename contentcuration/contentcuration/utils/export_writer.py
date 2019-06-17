# -*- coding: utf-8 -*-
import csv
import logging as logmodule
import math
import os
import sys
import tempfile
from collections import OrderedDict

import numpy as np
import pdfkit
from django.conf import settings
from django.contrib.sites.models import Site
from django.core.files.storage import default_storage
from django.template.loader import get_template
from django.utils.translation import ngettext
from django.utils.translation import ugettext as _
from le_utils.constants import content_kinds
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.text.fonts import FontFiles
from pptx.util import Inches
from pptx.util import Pt
from pressurecooker.encodings import encode_file_to_base64
from pressurecooker.encodings import write_base64_to_file
from wordcloud import WordCloud

from contentcuration.models import Channel
from contentcuration.models import ContentKind
from contentcuration.utils.files import generate_thumbnail_from_channel
from contentcuration.utils.format import format_size

# On OS X, the default backend will fail if you are not using a Framework build of Python,
# e.g. in a virtualenv. To avoid having to set MPLBACKEND each time we use Studio,
# automatically set the backend.
if sys.platform.startswith("darwin"):
    import matplotlib
    if matplotlib.get_backend().lower() == "macosx":
        matplotlib.use('PS')

import matplotlib.pyplot as plt  # noqa: E402


AUDIO_COLOR = "#F06292"
DOCUMENT_COLOR = "#FF3D00"
EXERCISE_COLOR = "#4DB6AC"
HTML_COLOR = "#FF8F00"
VIDEO_COLOR = "#283593"
SLIDESHOW_COLOR = "#4ECE90"

plt.switch_backend('agg')  # Avoid using tkinter as it causes server to stall (https://discuss.erpnext.com/t/wkhtmltopdf-error-erpnext-v7/14673/10)
os.environ['QT_QPA_PLATFORM'] = 'offscreen'  # Must be set for tests to run (https://github.com/ipython/ipython/issues/10627)
logmodule.basicConfig()
logging = logmodule.getLogger(__name__)


def _monkeypatch_font_directories():
    # python-pptx automatically fails on linux systems, so patch it
    # https://github.com/scanny/python-pptx/blob/master/pptx/text/fonts.py#L57
    def _can_i_haz_linux(cls):

        if sys.platform.startswith("linux"):
            return {
                # python-pptx fails if Calibri isn't found, so reroute it to a local font file
                ('Calibri', False, False): '/usr/share/fonts/truetype/freefont/FreeSans.ttf',
            }
        else:
            return FontFiles._old_installed_fonts()
    FontFiles._old_installed_fonts = FontFiles._installed_fonts
    FontFiles._installed_fonts = classmethod(_can_i_haz_linux)


_monkeypatch_font_directories()


class PDFMixin(object):
    def write_pdf(self, template, context, filepath, extra_options=None):
        template = get_template(template)
        html = template.render(context)
        options = {
            "encoding": "utf-8-sig",
            "quiet": "",
            'page-size': 'Letter',
            'margin-top': '0.5in',
            'margin-right': '0.5in',
            'margin-bottom': '0.5in',
            'margin-left': '0.5in',
        }
        if extra_options:
            options.update(extra_options)
        pdfkit.from_string(html, filepath, options=options)
        return filepath


class PPTMixin(object):
    slide = None
    width = float(10)
    height = float(5.6)

    def get_next_slide(self):
        self.ppt.slide_width = Inches(self.width)
        self.ppt.slide_height = Inches(self.height)
        slide_layout = self.ppt.slide_layouts[6]  # Get a blank slide
        self.slide = self.ppt.slides.add_slide(slide_layout)

    def get_rgb_from_hex(self, hexval):
        hexval = hexval.upper()
        mapping = {"0": 0, "1": 1, "2": 2, "3": 3, "4": 4, "5": 5, "6": 6, "7": 7, "8": 8, "9": 9, "A": 10, "B": 11, "C": 12, "D": 13, "E": 14, "F": 15}
        red = mapping[hexval[1]] * 16 + mapping[hexval[2]]
        green = mapping[hexval[3]] * 16 + mapping[hexval[4]]
        blue = mapping[hexval[5]] * 16 + mapping[hexval[6]]
        return RGBColor(red, green, blue)

    def generate_textbox(self, left=0, top=0, width=12, height=1, word_wrap=True):
        left = Inches(left)
        top = Inches(top)
        width = Inches(width)
        height = Inches(height)
        textbox = self.slide.shapes.add_textbox(left, top, width, height)
        textframe = textbox.text_frame
        textframe.word_wrap = word_wrap
        return textframe

    def add_line(self, textframe, text, fontsize=12, bold=False, color=None, italic=False, space_after=0.5, space_before=0, append=True):
        p = textframe.add_paragraph() if append else textframe.paragraphs[0]
        p.space_after = Pt(space_after)
        p.space_before = Pt(space_before)
        self.add_run(p, text, fontsize=fontsize, bold=bold, color=color, italic=italic)
        return p

    def add_run(self, paragraph, text, fontsize=12, bold=False, color=None, italic=False):
        run = paragraph.add_run()
        run.font.size = Pt(fontsize)
        run.font.name = 'Calibri'
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color or RGBColor(0, 0, 0)
        run.text = text
        return run

    def get_thumbnail_from_encoding(self, encoding):
        filepath = self.get_write_to_path(ext="png")
        write_base64_to_file(encoding, filepath)
        return filepath

    def add_picture(self, encoding, left=0, top=0, width=2, height=2):
        filepath = self.get_write_to_path(ext="png")
        write_base64_to_file(encoding, filepath)
        return self.slide.shapes.add_picture(filepath, Inches(left), Inches(top), width=Inches(width), height=Inches(height))

    def add_shape(self, shape=MSO_SHAPE.RECTANGLE, left=0, top=0, width=1, height=1, color=None):
        shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color or RGBColor(0, 0, 0)
        shape.line.color.rgb = color or RGBColor(0, 0, 0)
        shape.shadow.inherit = False
        return shape


class CSVMixin(object):
    def write_csv(self, filepath, rows, header=None):
        with open(filepath, 'wb') as csvfile:
            writer = csv.writer(csvfile, delimiter=',', quoting=csv.QUOTE_MINIMAL)
            if header:
                writer.writerow(header)
            for row in rows:
                writer.writerow(row)

        return filepath


class ExportWriter(object):
    tempfiles = None
    ext = None

    def __init__(self, *args, **kwargs):
        self.tempfiles = []

    def pluralize_constant(self, count, constant):
        data = {'count': count}
        if constant == content_kinds.TOPIC:
            return ngettext('%(count)d Topic', '%(count)d Topics', count) % data
        elif constant == content_kinds.VIDEO:
            return ngettext('%(count)d Video', '%(count)d Videos', count) % data
        elif constant == content_kinds.AUDIO:
            return ngettext('%(count)d Audio', '%(count)d Audios', count) % data
        elif constant == content_kinds.EXERCISE:
            return ngettext('%(count)d Exercise', '%(count)d Exercises', count) % data
        elif constant == content_kinds.DOCUMENT:
            return ngettext('%(count)d Document', '%(count)d Documents', count) % data
        elif constant == content_kinds.HTML5:
            return ngettext('%(count)d Html App', '%(count)d Html Apps', count) % data
        elif constant == content_kinds.SLIDESHOW:
            return ngettext('%(count)d Slideshow', '%(count)d Slideshows', count) % data
        elif constant == "resource":
            return ngettext('%(count)d Total Resource', '%(count)d Total Resources', count) % data
        elif constant == "resource_split":
            return ngettext('%(count)d\nTotal Resource', '%(count)d\nTotal Resources', count) % data
        else:
            logging.warning('No translation for pluralizing {}'.format(constant))
            return '{} {}'.format(count, constant)

    def get_write_to_path(self, ext=None):
        ext = ext or self.ext
        tempf = tempfile.NamedTemporaryFile(suffix=".{}".format(ext), delete=False)
        self.tempfiles.append(tempf.name)
        return tempf.name

    def write(self, *args, **kwargs):
        raise NotImplementedError("Must implement a write method for this class")

    def delete_tempfiles(self):
        for tempf in self.tempfiles:
            os.unlink(tempf)
        self.tempfiles = []


class ChannelDetailsWriter(ExportWriter):
    # Needs to be alphabetized to match content kind sorting
    color_selection = [AUDIO_COLOR, DOCUMENT_COLOR, EXERCISE_COLOR, HTML_COLOR, SLIDESHOW_COLOR, VIDEO_COLOR]
    condensed_tag_limit = 10
    size_divisor = 100000000
    scale_text = [_("Very Small")] * 2 + [_("Small")] * 2 + [_("Average")] * 3 + [_("Large")] * 2 + [_("Very Large")] * 2
    tagcloud_width = 600
    tagcloud_height = None

    def __init__(self, channel_ids, site=None, condensed=False, filename=None):
        super(ChannelDetailsWriter, self).__init__()
        self.channels = Channel.objects.filter(pk__in=channel_ids)  # Implementing as a list so we can easily make this apply to bundles
        self.filename = filename
        if self.channels.count() == 1 and not filename:
            self.filename = self.channels[0].pk
        elif not filename:
            raise ValueError("Must specify a filename if channel count is greater than 1")
        self.site = site or Site.objects.get(id=1)
        self.condensed = condensed

    def write(self, *args, **kwargs):
        try:
            filepath = self.get_write_to_path()
            self._write_details(filepath)

            saved_filename = "{}.{}".format(self.filename, self.ext)
            save_to_path = os.path.sep.join([settings.EXPORT_ROOT, saved_filename])

            # Write file to default storage
            with open(filepath, 'rb') as fobj:
                default_storage.save(save_to_path, fobj)
            return save_to_path

        finally:
            self.delete_tempfiles()

    def _write_details(self, *args, **kwargs):
        raise NotImplementedError("Must implement a write_export_file method for ChannelDetailsWriter subclasses")

    def get_channel_data(self, channel):
        data = channel.main_tree.get_details()
        primarytoken = channel.secret_tokens.filter(is_primary=True).first()

        data.update({
            "channel": channel,
            "site": 'https://' + self.site.domain,
            "thumbnail": generate_thumbnail_from_channel(channel, dimension=300) or self.get_default_thumbnail_encoding(),
            "tokens": [str(t) for t in channel.secret_tokens.exclude(token=channel.pk).filter(is_primary=True)],
            "primarytoken": primarytoken and str(primarytoken),
            "storage": self.get_storage_bar(data['resource_size']),
            "size": self.get_size_bar(data['resource_count']),
            "piechart": self.get_pie_chart(data['kind_count'], small_layout=self.condensed),
            "tagcloud": data['tags'] and self.get_tagcloud(data['tags'], tag_limit=self.condensed and self.condensed_tag_limit),
        })
        return data

    def get_default_thumbnail_encoding(self):
        try:
            filepath = os.path.join(settings.STATIC_ROOT, 'img', 'kolibri_placeholder.png')
            return encode_file_to_base64(filepath, "data:image/png;base64,")
        except IOError:
            logging.warning("Could not find {}".format(filepath))

    def get_storage_bar(self, size):
        try:
            size_index = int(max(1, min(math.ceil(math.log(size/self.size_divisor, 2)), 10)))
        except ValueError:
            size_index = 1
        return {
            "filled": range(size_index),
            "text": self.scale_text[size_index],
            "storage": "{} {}".format(*format_size(size)),
        }

    def get_size_bar(self, count):
        try:
            size_index = int(max(1, min(math.floor(math.log(count, 2.8)), 10)))
        except ValueError:
            size_index = 1

        return {
            "filled": size_index,
            "scale": range(len(self.scale_text)),
            "text": self.scale_text[size_index]
        }

    def get_pie_chart(self, counts, small_layout=False):
        # Put kind counts in a usable format
        kinds = list(ContentKind.objects.exclude(kind=content_kinds.TOPIC)
                                .order_by('kind')
                                .values_list('kind', flat=True))
        kind_vals = {k: next((c['count'] for c in counts if c['kind_id'] == k), 0) for k in kinds}
        kind_vals = OrderedDict(sorted(kind_vals.items()))
        sizes = [v for k, v in kind_vals.items()]
        total = max(sum(sizes), 1)

        labels = [{
            "text": ' {text} \n{p:.1f}%'.format(
                    text=self.pluralize_constant(v, k),
                    p=float(v)/total * 100.0
            ),
            "count": v
        } for k, v in kind_vals.items()]

        # Create pie chart
        fig, ax = plt.subplots(subplot_kw=dict(aspect="equal"))
        wedgeprops = {"edgecolor": "white", 'linewidth': 1, 'linestyle': 'solid', 'antialiased': True}
        wedges, texts = ax.pie(sizes, colors=self.color_selection, wedgeprops=wedgeprops)

        bbox_props = dict(boxstyle="square,pad=0.3", fc="w", ec="k", lw=0.72)
        kw = dict(xycoords='data', textcoords='data', arrowprops=dict(arrowstyle="-"),
                  bbox=bbox_props, zorder=0, va="center")

        # Add popout labels for the larger layout
        if not small_layout:
            for i, p in enumerate(wedges):
                if not labels[i]['count']:
                    continue

                ang = (p.theta2 - p.theta1)/2. + p.theta1
                y = np.sin(np.deg2rad(ang))
                x = np.cos(np.deg2rad(ang))
                connectionstyle = "angle,angleA=0,angleB={}".format(ang)
                kw["arrowprops"].update({"connectionstyle": connectionstyle, "facecolor": "gray"})
                ax.annotate(labels[i]['text'], xy=(x, y), xytext=(1.35*np.sign(x), 1.4*y),
                            ha="center", fontsize=10, **kw)

        # Add legend for the smaller layout
        else:
            plt.legend(
                loc='center left',
                labels=[l['text'].split('\n')[0] for l in labels],
                prop={'size': 20},
                bbox_to_anchor=(0.7, 0.5),
                bbox_transform=plt.gcf().transFigure
            )

        # Set up size variables for center circle
        center_text_size = 25 if small_layout else 20  # Renders smaller, so text needs to be bigger
        center_text_ratio = 0.75 if small_layout else 0.6

        # Add center circle
        circle = plt.Circle((0, 0), center_text_ratio, fc='white')
        centertext = self.pluralize_constant(sum(sizes), "resource_split").split("\n")
        plt.annotate(centertext[0], xy=(0, 0.1), fontsize=center_text_size, ha="center")
        plt.annotate(centertext[1], xy=(0, -0.15), fontsize=center_text_size - 5, ha="center")
        fig = plt.gcf()
        fig.gca().add_artist(circle)
        plt.tight_layout()

        # Write chart to image and get encoding
        filepath = self.get_write_to_path(ext="png")
        plt.savefig(filepath, bbox_inches='tight')
        plt.clf()
        plt.close()
        return encode_file_to_base64(filepath, "data:image/png;base64,")

    def get_tagcloud(self, tags, tag_limit=None):
        tag_limit = tag_limit or len(tags)
        tags = sorted(tags, key=lambda kv: -kv['count'])[:tag_limit]  # Get top X tags
        tag_dict = {t['tag_name']: t['count'] for t in tags}

        # Generate a word cloud image
        wordcloud = WordCloud(
            background_color='white',
            min_font_size=10,
            max_font_size=60,
            width=self.tagcloud_width,
            height=self.tagcloud_height or 30 * len(tags) / 2 + 10,
            font_path=os.path.sep.join([settings.STATIC_ROOT, 'fonts', 'OpenSans-Regular.ttf'])
        ).generate_from_frequencies(tag_dict)

        tag_counts = [t['count'] for t in tags]
        step = (float(max(tag_counts))) / len(self.color_selection)
        thresholds = list(reversed([int(round(i * step)) for i in range(len(self.color_selection))]))

        def get_color(word, font_size, position, orientation, random_state=None, **kwargs):
            index = next((i for i, t in enumerate(thresholds) if tag_dict[word] >= t), 0)
            return self.color_selection[index]

        wordcloud.recolor(color_func=get_color)
        image = wordcloud.to_image()
        filepath = self.get_write_to_path(ext="png")
        image.save(filepath)
        return encode_file_to_base64(filepath, "data:image/png;base64,")


class ChannelDetailsPDFWriter(ChannelDetailsWriter, PDFMixin):
    ext = "pdf"

    def __init__(self, channel_ids, condensed=False, **kwargs):
        super(ChannelDetailsPDFWriter, self).__init__(channel_ids, condensed=condensed, **kwargs)
        self.filename = "{} (condensed)".format(self.filename) if condensed else self.filename
        self.template = "export/channel_detail_pdf_condensed.html" if condensed else "export/channel_detail_pdf.html"

    def _write_details(self, filepath):
        if self.channels.count() == 1:
            footer_text = _("Page %(page)s of %(pagecount)s - %(channel)s can be found on Kolibri Studio, a product of Learning Equality") \
                            % {"page": "[page]", "pagecount": "[topage]", "channel": self.channels[0].name[:40]}
        else:
            footer_text = _("Page %(page)s of %(pagecount)s - These channels can be found on Kolibri Studio, a product of Learning Equality") \
                            % {"page": "[page]", "pagecount": "[topage]"}

        data = {
            "channels": [self.get_channel_data(channel) for channel in self.channels],
            "colors": {
                "audio": AUDIO_COLOR,
                "document": DOCUMENT_COLOR,
                "exercise": EXERCISE_COLOR,
                "html": HTML_COLOR,
                "video": VIDEO_COLOR,
                "slideshow": SLIDESHOW_COLOR,
            }
        }
        try:
            self.write_pdf(self.template, data, filepath, extra_options={"footer-center": footer_text, "footer-font-size": "9"})
        except IOError as e:
            logging.error("Unable to generate PDF, attempting without footer: {}".format(str(e)))
            self.write_pdf(self.template, data, filepath)


class ChannelDetailsPPTWriter(ChannelDetailsWriter, PPTMixin):
    ext = "pptx"
    tagcloud_width = 430
    tagcloud_height = 210
    condensed_tag_limit = 20
    gray = RGBColor(170, 170, 170)

    def __init__(self, channel_ids, **kwargs):
        super(ChannelDetailsPPTWriter, self).__init__(channel_ids, condensed=True, **kwargs)

    def _write_details(self, filepath):
        self.ppt = Presentation()
        for channel in self.channels:
            self.write_slide(channel)

        # Save the file
        self.ppt.save(filepath)

    def write_slide(self, channel):
        self.get_next_slide()
        data = self.get_channel_data(channel)

        next_line = 0.1  # Keeps track of last line (useful for setting the top location of shapes)

        # Add thumbnail
        padding = 0.2
        thumbnail_width = 1.1
        if data['thumbnail']:
            thumbnail = self.add_picture(data['thumbnail'], padding, padding, thumbnail_width, thumbnail_width)
            thumbnail.line.color.rgb = self.gray
            thumbnail.line.width = Inches(0.01)

        # Add title/description
        title_left = thumbnail_width + padding * 2

        title_height = 0.5
        title_tf = self.generate_textbox(title_left,  next_line,  self.width - title_left, title_height)
        title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        self.add_line(title_tf, channel.name, fontsize=24, bold=True, append=False)
        next_line += title_height

        # Variables for section under title
        includes_width = 2
        size_height = 0.5
        description_height = 1.25
        size_width = self.width - title_left - includes_width - padding * 2

        # Add language information
        icon_width = 0.2
        language_left = size_width + title_left + padding
        language_icon_path = os.path.join(settings.STATIC_ROOT, 'img', 'export', 'language.png')
        encoding = encode_file_to_base64(language_icon_path, 'data:image/png;base64,')
        self.add_picture(encoding, language_left, next_line + 0.04, icon_width, icon_width)

        includes_tf = self.generate_textbox(language_left + icon_width - 0.08,  next_line,  includes_width, size_height + description_height)
        language = channel.language.native_name if channel.language else _("No language set")
        self.add_line(includes_tf, " {}".format(language), append=False, bold=True)
        if data['accessible_languages']:
            self.add_line(includes_tf, _("    * Subtitles included"), fontsize=10, space_before=2)

        # Add For Educators: Coach Content
        if data['includes'].get('coach_content'):
            coach_content = self.add_line(includes_tf, "✔", bold=True, color=self.get_rgb_from_hex(EXERCISE_COLOR), space_before=4)
            self.add_run(coach_content, _(" Coach Content"))

        # Add For Educators: Assessments
        if data['includes'].get('exercises'):
            assessments = self.add_line(includes_tf, "✔", bold=True, color=self.get_rgb_from_hex(EXERCISE_COLOR), space_before=4)
            self.add_run(assessments, _(" Assessments"))

        # Add size information
        size_tf = self.generate_textbox(title_left,  next_line,  size_width, size_height)
        size_bar = self.add_line(size_tf, "", append=False)
        for i in data['size']['scale']:
            self.add_run(size_bar, "▮", color=self.get_rgb_from_hex(EXERCISE_COLOR) if i < data['size']['filled'] else self.gray, fontsize=14)
        self.add_line(size_tf, _("Channel size: %(size)s") % {"size": data['size']['text'].lower()}, fontsize=8, italic=True, color=self.gray)
        next_line += size_height

        # Add description
        description_tf = self.generate_textbox(title_left,  next_line, size_width, description_height)
        self.add_line(description_tf, channel.description, color=self.gray, append=False)
        description_tf.fit_text()
        next_line += description_height + 0.1

        # Add separator with headers
        separator_height = 0.3
        self.add_shape(left=0, top=next_line, width=self.width/2, height=separator_height, color=self.get_rgb_from_hex(EXERCISE_COLOR))
        resource_header = self.generate_textbox(padding, next_line, self.width / 2 - padding, separator_height)
        self.add_line(resource_header, _("Resource Breakdown"), bold=True, color=self.get_rgb_from_hex("#FFFFFF"), append=False)

        self.add_shape(left=self.width/2, top=next_line, width=self.width/2, height=separator_height, color=self.get_rgb_from_hex("#595959"))
        tag_header = self.generate_textbox(padding + self.width / 2 - padding, next_line, self.width / 2 - padding, separator_height)
        self.add_line(tag_header, _("Most Common Tags"), bold=True, color=self.get_rgb_from_hex("#FFFFFF"), append=False)
        next_line += separator_height + 0.05

        # Add piechart
        chart_height = 2.3
        if data['resource_count']:
            self.add_picture(data['piechart'], 0, next_line, self.width / 2 - 1, height=chart_height)
        else:
            empty_tf = self.generate_textbox(0,  next_line, self.width / 2, chart_height)
            empty_line = self.add_line(empty_tf, _("No Resources Found"), color=self.gray, fontsize=14, italic=True)
            empty_line.alignment = PP_ALIGN.CENTER

        # Add tagcloud
        if data['tags']:
            self.add_picture(data['tagcloud'], self.width/2 + padding, next_line + 0.1, self.width / 2 - 1, chart_height - padding * 2)
        else:
            empty_tf = self.generate_textbox(self.width / 2,  next_line, self.width / 2, chart_height)
            empty_line = self.add_line(empty_tf, _("No Tags Found"), color=self.gray, fontsize=14, italic=True)
            empty_line.alignment = PP_ALIGN.CENTER
        next_line += chart_height + 0.01

        # Add logo
        logo_width = 0.9
        logo_height = 0.25
        logo_left = Inches(self.width / 2 - logo_width / 2)
        try:
            logo_url = os.path.join(settings.STATIC_ROOT, 'img', 'le_login.png')
            self.slide.shapes.add_picture(logo_url, logo_left, Inches(next_line), width=Inches(logo_width), height=Inches(logo_height))
        except IOError:
            logging.warning("Unable to add LE logo")
        next_line += logo_height

        # Add disclaimer
        disclaimer_tf = self.generate_textbox(0, next_line,  self.width, 0.2)
        disclaimer_line = self.add_line(disclaimer_tf, _("This slide was automatically generated by Kolibri Studio, a product of Learning Equality"),
                                        fontsize=7, color=self.gray, append=False)
        disclaimer_line.alignment = PP_ALIGN.CENTER


class ChannelDetailsCSVWriter(ChannelDetailsWriter, CSVMixin):
    ext = "csv"

    def _write_details(self, filepath):
        header = [_("Name"), _("Description"), _("Language"), _("Token"), _("Size"), _("Storage"), _("Resources"),
                  _("Languages"), _("Subtitles"), _("Coach Content?"), _("Assessments?"), _("Tags"), _("Authors"),
                  _("Providers"), _("Aggregators"), _("Licenses"), _("Copyright Holders")]
        rows = []

        for channel in self.channels:
            data = self.get_channel_data(channel)
            language = channel.language.native_name if channel.language else _("No language set")
            token = data['primarytoken'] if data['primarytoken'] else _("Publish channel to get token")
            size = "{} - {}".format(self.pluralize_constant(data['resource_count'], "resource"), data['size']['text'])
            storage = "{} - {}".format(data['storage']['storage'], data['storage']['text'])
            resources = " | ".join([self.pluralize_constant(k['count'], k['kind_id']) for k in data['kind_count']])
            languages = " | ".join(data['languages'])
            subtitles = " | ".join(data['accessible_languages'])
            coach_content = _("Yes") if data['includes']['coach_content'] else _("No")
            assessments = _("Yes") if data['includes']['exercises'] else _("No")
            tags = " | ".join([t['tag_name'] for t in data['tags']])
            authors = " | ".join(data['authors'])
            providers = " | ".join(data['providers'])
            aggregators = " | ".join(data['aggregators'])
            licenses = " | ".join(data['licenses']).encode('utf-8')
            copyright_holders = " | ".join(data['copyright_holders'])

            rows.append([channel.name, channel.description, language, token, size, storage, resources,
                         languages, subtitles, coach_content, assessments, tags, authors, providers,
                         aggregators, licenses, copyright_holders])

        return self.write_csv(filepath, rows, header=header)
